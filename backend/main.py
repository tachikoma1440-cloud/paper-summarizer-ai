import asyncio
import base64
import datetime
import io
import math
import os
import pathlib
import re
from urllib.parse import quote

import fitz  # PyMuPDF
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
import anthropic
from PIL import Image as PILImage
from pydantic import BaseModel
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

app = FastAPI(title="Paper Summarizer API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

client       = anthropic.Anthropic(api_key=os.environ.get("ANTHROPIC_API_KEY"))
async_client = anthropic.AsyncAnthropic(api_key=os.environ.get("ANTHROPIC_API_KEY"))

OBSIDIAN_VAULT_PATH = os.environ.get("OBSIDIAN_VAULT_PATH", "")
DEV_LOG_PATH        = pathlib.Path("/app/dev_log.md")

# ── Pricing constants (Claude claude-sonnet-4-6) ──────────────────────────────
PRICE_INPUT_TEXT_USD  = 3.00    # USD per 1M input tokens
PRICE_OUTPUT_USD      = 15.00   # USD per 1M output tokens
TOKENS_PER_IMAGE      = 1_600   # approximate input tokens per image
USD_TO_JPY            = 150     # exchange rate (conservative, round up)
MAX_TOKENS_SUMMARIZE  = 4_096   # max_tokens used in /summarize
MAX_TOKENS_FIGURE     = 500     # max_tokens used in _analyze_figure_async

_BLACK = RGBColor(0x00, 0x00, 0x00)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_NAVY  = RGBColor(0x1F, 0x49, 0x7D)
_LGRAY = RGBColor(0xF2, 0xF2, 0xF2)
_MGRAY = RGBColor(0x60, 0x60, 0x60)

# ── Prompts ───────────────────────────────────────────────────────────────────

SYSTEM_PROMPT = """あなたは医学論文の抄読会資料を作成する専門家です。
論文のIMRAD構造（Introduction, Methods, Results, Discussion）を正確に把握したうえで、
指定された5項目を日本語のMarkdown形式で出力してください。

出力ルール：
- 各項目は ## 見出しで始める
- 数値・統計は原文のまま記載する（改変しない）
- 情報が論文中に明示されていない場合は「記載なし」と書く
- 推測や補足は加えない
- 主要な疾患名・薬剤名・統計手法・医学的概念は [[用語]] の形式で囲む（例：[[心筋梗塞]]、[[ランダム化比較試験]]、[[カプラン・マイヤー法]]、[[Cox比例ハザードモデル]]）"""


def build_prompt(paper_text: str) -> str:
    return f"""以下の医学論文テキストを読み、次の5項目をMarkdown形式で出力してください。

---

## 1. PICO

| 項目 | 内容 |
|------|------|
| **P** (Patient/Population) | 対象患者・集団 |
| **I** (Intervention) | 介入・曝露 |
| **C** (Comparison) | 比較対照 |
| **O** (Outcome) | 主要アウトカム |

上記の表を論文の内容で埋めてください。

---

## 2. この研究の新規性と重要性

（既存研究との差別化ポイント、なぜ今この研究が必要か）

---

## 3. 主要な結果（数値を含む）

（主要アウトカムの結果を、p値・信頼区間・効果量などの統計値とともに記載）

---

## 4. この論文の限界（Limitations）

（著者が認めている限界点を箇条書きで列挙）

---

## 5. 臨床現場でどう活かせるか（Clinical Implications）

（実際の診療・ケアへの応用可能性）

---

論文テキスト：
{paper_text[:15000]}
"""

# ── Progress report ───────────────────────────────────────────────────────────

PROGRESS_SYSTEM_PROMPT = """あなたはプロジェクト進捗報告資料を作成する専門家です。
開発ログを分析し、指定された4項目を日本語のMarkdown形式で出力してください。
各項目は ## 見出しで始め、内容は箇条書きで記載してください。"""


def build_progress_prompt(dev_log: str) -> str:
    return f"""以下の開発ログを読み、プロジェクト進捗報告スライド用の情報を以下4項目のMarkdown形式で出力してください。

## 開発タイムライン
（日付ごとの主要作業を時系列で箇条書き）

## 実装した主要機能
（実装済みの機能・改善点を箇条書き）

## 課題と解決策
（開発中の課題と解決策。記録が少ない場合は技術的チャレンジを推察して記載）

## 今後の展望・改善点
（改善案・追加機能の候補を箇条書き）

開発ログ：
{dev_log[:8000]}
"""


FIXED_PROJECT_SECTIONS: list[dict] = [
    {
        "title": "プロジェクト概要",
        "content": (
            "- 医学論文のPDFをアップロードするだけで抄読会資料を自動生成するAIアプリ\n"
            "- PICO・新規性・結果・限界・臨床応用の5項目を日本語で自動抽出\n"
            "- Claude API（claude-sonnet-4-6）によるマルチモーダル解析で図表も自動識別・解説\n"
            "- MarkdownおよびPowerPointスライドとして即時ダウンロード可能"
        ),
    },
    {
        "title": "システム構成",
        "content": (
            "- フロントエンド：Streamlit（Pythonベースの直感的UI）\n"
            "- バックエンド：FastAPI（高速REST APIサーバー）\n"
            "- AI：Anthropic Claude API（claude-sonnet-4-6）\n"
            "- 図表処理：PyMuPDF（PDF解析） + Pillow（画像処理）\n"
            "- スライド生成：python-pptx\n"
            "- インフラ：Docker Compose（フロントエンド・バックエンド 2コンテナ構成）"
        ),
    },
    {
        "title": "主要機能一覧",
        "content": (
            "- 論文PDF → 5項目サマリー自動生成（PICO・新規性・結果・限界・臨床応用）\n"
            "- Markdownファイルとしてダウンロード\n"
            "- PowerPointスライド自動生成（テキスト版）\n"
            "- 図表入りPowerPointスライド生成（AIによる図表解説付き）\n"
            "- Obsidian Vault連携（論文要約・開発ログの外部記憶化）\n"
            "- プロジェクト進捗報告PPTX自動生成"
        ),
    },
]

# ── PDF utilities ─────────────────────────────────────────────────────────────

def extract_text(pdf_bytes: bytes) -> str:
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    return "\n".join(page.get_text() for page in doc)


# ── Image utilities ───────────────────────────────────────────────────────────

def _autocrop_white(img: PILImage.Image, threshold: int = 248, padding: int = 8) -> PILImage.Image:
    """Crop near-white margins to maximize figure area on slide."""
    mask = img.convert("L").point(lambda p: 255 if p < threshold else 0)
    bbox = mask.getbbox()
    if not bbox:
        return img
    w, h = img.size
    return img.crop((
        max(0, bbox[0] - padding),
        max(0, bbox[1] - padding),
        min(w, bbox[2] + padding),
        min(h, bbox[3] + padding),
    ))


def _is_real_figure(img: PILImage.Image) -> bool:
    """
    Return True if the image looks like a graph / photo / table.
    Reject text-heavy document fragments (mostly white + black with no midtones/color).
    """
    w, h = img.size

    # Very tall-and-narrow strips are usually column fragments, not figures
    if w / h < 0.25:
        return False

    # Use an 80×80 thumbnail for fast pixel analysis
    thumb = img.convert("RGB").resize((80, 80), PILImage.LANCZOS)
    gray  = thumb.convert("L")
    hist  = gray.histogram()          # 256 bins, total = 6400
    total = 80 * 80

    white_ratio = sum(hist[230:]) / total   # near-white  (230-255)
    dark_ratio  = sum(hist[:25])  / total   # near-black  (0-24)
    mid_ratio   = sum(hist[25:230]) / total  # midtones

    # Signature of a text/document image: lots of white, some black, barely any midtone
    if white_ratio > 0.82 and dark_ratio > 0.01 and mid_ratio < 0.18:
        # Allow if the image has significant colour (e.g. colour bar charts)
        def _wmean(ch_hist):
            return sum(i * c for i, c in enumerate(ch_hist)) / total

        r_m = _wmean(thumb.getchannel("R").histogram())
        g_m = _wmean(thumb.getchannel("G").histogram())
        b_m = _wmean(thumb.getchannel("B").histogram())
        color_spread = max(abs(r_m - g_m), abs(g_m - b_m), abs(r_m - b_m))

        if color_spread < 12:   # Essentially greyscale → text, reject
            return False

    return True


def _load_and_normalize(raw: bytes, max_dim: int = 1600) -> PILImage.Image | None:
    """Open raw image bytes, convert to RGB, downscale if too large."""
    try:
        img = PILImage.open(io.BytesIO(raw))
        if max(img.width, img.height) > max_dim:
            ratio = max_dim / max(img.width, img.height)
            img = img.resize(
                (int(img.width * ratio), int(img.height * ratio)),
                PILImage.LANCZOS,
            )
        if img.mode == "P":
            img = img.convert("RGBA")
        if img.mode in ("RGBA", "LA"):
            bg = PILImage.new("RGB", img.size, (255, 255, 255))
            bg.paste(img, mask=img.split()[-1])
            return bg
        if img.mode != "RGB":
            return img.convert("RGB")
        return img
    except Exception:
        return None


def _to_png_bytes(img: PILImage.Image) -> bytes:
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


_CAPTION_RE = re.compile(
    r'^(Fig\.?|Figure|Table|Tab\.?|FIG\.?|TABLE|FIGURE|Suppl\.?\s*Fig)',
    re.IGNORECASE,
)


def _find_caption(page, xref: int, text_blocks: list) -> str:
    """Return the Figure/Table legend nearest to this image on the page."""
    try:
        rects = page.get_image_rects(xref)
        if not rects:
            return ""
        ir = rects[0]   # image rect on the page

        best, best_dist = "", float("inf")
        for blk in text_blocks:
            if len(blk) < 5:
                continue
            bx0, by0, bx1, by1, text = blk[0], blk[1], blk[2], blk[3], blk[4]
            if not isinstance(text, str) or not text.strip():
                continue
            text = text.strip()
            if not _CAPTION_RE.match(text):
                continue

            dist_below = by0 - ir.y1   # positive → block is below image
            dist_above = ir.y0 - by1   # positive → block is above image

            in_below = -5 <= dist_below <= 120
            in_above = -5 <= dist_above <= 80
            if not (in_below or in_above):
                continue

            # Require horizontal overlap or near-alignment
            h_overlap = min(float(bx1), float(ir.x1)) - max(float(bx0), float(ir.x0))
            if h_overlap < -30:
                continue

            # When both directions qualify, pick the smaller absolute gap.
            # Use abs() for comparison so slight overlaps (negative) never
            # beat a clean positive gap.
            if in_below and in_above:
                dist = dist_below if abs(dist_below) <= abs(dist_above) else dist_above
            elif in_below:
                dist = dist_below
            else:
                dist = dist_above

            if abs(dist) < best_dist:
                best_dist = abs(dist)
                best = text[:400]

        return best
    except Exception:
        return ""


# ── Screenshot rendering helpers ──────────────────────────────────────────────

def _render_region(page, rect, padding_pt: float = 10) -> bytes:
    """Render a rectangular region of a PDF page to PNG at 2× resolution."""
    pr   = page.rect
    clip = fitz.Rect(
        max(0, rect.x0 - padding_pt),
        max(0, rect.y0 - padding_pt),
        min(pr.width,  rect.x1 + padding_pt),
        min(pr.height, rect.y1 + padding_pt),
    )
    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), clip=clip)
    return pix.tobytes("png")


def _find_title_cut_y(page) -> float:
    """
    Return the y-coordinate where the abstract / body text begins on page 1.
    Everything above this line is considered title + authors + institution.
    """
    BODY_KW = {
        'abstract', 'introduction', 'background', 'methods', 'keyword',
        'objective', 'purpose', 'summary', '要旨', '抄録', '背景',
    }
    blocks = sorted(page.get_text("blocks"), key=lambda b: b[1])
    page_h = page.rect.height
    for blk in blocks:
        if len(blk) < 5:
            continue
        y0, text = blk[1], blk[4].strip()
        if len(text) < 2 or y0 < page_h * 0.12:
            continue
        if len(text) <= 100 and any(kw in text.lower() for kw in BODY_KW):
            return max(y0 - 8, page_h * 0.15)
    return page_h * 0.42   # fallback: top 42 %


def _expand_for_caption(img_rect, caption: str, text_blocks: list, page_rect):
    """Expand img_rect to encompass the matching caption text block."""
    if not caption:
        return img_rect
    sample = caption[:60]
    for blk in text_blocks:
        if len(blk) < 5:
            continue
        bx0, by0, bx1, by1, text = blk[0], blk[1], blk[2], blk[3], blk[4]
        if sample not in text and text[:60] not in caption:
            continue
        return fitz.Rect(
            min(img_rect.x0, float(bx0)),
            min(img_rect.y0, float(by0)),
            max(img_rect.x1, float(bx1)),
            max(img_rect.y1, float(by1)),
        )
    return img_rect


_TABLE_CAPTION_RE = re.compile(
    r'^(Table|Tab\.?|TABLE|Suppl\.?\s*Table)',
    re.IGNORECASE,
)


def _find_table_caption(tbl_bbox, text_blocks: list) -> str:
    """Return the nearest Table X. caption for a detected table bounding box."""
    tr = fitz.Rect(tbl_bbox)
    best, best_dist = "", float("inf")
    for blk in text_blocks:
        if len(blk) < 5:
            continue
        bx0, by0, bx1, by1, text = blk[0], blk[1], blk[2], blk[3], blk[4]
        text = text.strip()
        if not _TABLE_CAPTION_RE.match(text):
            continue
        dist_above = tr.y0 - float(by1)
        dist_below = float(by0) - tr.y1
        if -5 <= dist_above <= 80 or -5 <= dist_below <= 120:
            dist = min(abs(dist_above), abs(dist_below))
            if dist < best_dist:
                best_dist = dist
                best = text[:400]
    return best


# ── Visual element extraction ─────────────────────────────────────────────────

def extract_figures(pdf_bytes: bytes, max_figures: int = 10) -> list[dict]:
    """
    Extract figures and tables as rendered page-region screenshots.
    Body text is never included — only the figure/table area + its caption.
    """
    doc   = fitz.open(stream=pdf_bytes, filetype="pdf")
    items: list[dict] = []
    seen:  set[int]   = set()

    for page_num, page in enumerate(doc, 1):
        text_blocks = page.get_text("blocks")
        page_rect   = page.rect
        page_area   = page_rect.width * page_rect.height

        # ── ① Figures: embedded image xrefs → render page clip ───────────
        for img_info in page.get_images(full=True):
            xref = img_info[0]
            if xref in seen:
                continue
            seen.add(xref)
            try:
                base = doc.extract_image(xref)
                w, h = base["width"], base["height"]
                if min(w, h) < 80 or w * h < 10_000:
                    continue

                img_rects = page.get_image_rects(xref)
                if not img_rects:
                    continue
                ir = img_rects[0]

                coverage = (ir.width * ir.height) / page_area
                if coverage > 0.55 or coverage < 0.005:
                    continue
                if ir.height > 0 and (ir.width / ir.height) < 0.25:
                    continue

                img_pil = _load_and_normalize(base["image"])
                if img_pil is None or not _is_real_figure(img_pil):
                    continue

                caption     = _find_caption(page, xref, text_blocks)
                render_rect = _expand_for_caption(ir, caption, text_blocks, page_rect)
                items.append({
                    "page_num":   page_num,
                    "figure_num": len(items) + 1,
                    "png_bytes":  _render_region(page, render_rect),
                    "caption":    caption,
                })
                if len(items) >= max_figures:
                    return items
            except Exception:
                continue

        # ── ② Tables: find_tables() → render page clip ───────────────────
        try:
            for tbl in page.find_tables().tables:
                caption     = _find_table_caption(tbl.bbox, text_blocks)
                tbl_rect    = fitz.Rect(tbl.bbox)
                render_rect = _expand_for_caption(tbl_rect, caption, text_blocks, page_rect)
                items.append({
                    "page_num":   page_num,
                    "figure_num": len(items) + 1,
                    "png_bytes":  _render_region(page, render_rect),
                    "caption":    caption,
                })
                if len(items) >= max_figures:
                    return items
        except Exception:
            pass   # find_tables() unavailable or no tables on this page

    return items


def extract_title_screenshot(pdf_bytes: bytes) -> bytes | None:
    """
    Render the title / authors / institution area of page 1 as a PNG screenshot.
    Crops just above where the abstract / body text begins, so body text is excluded.
    """
    try:
        doc  = fitz.open(stream=pdf_bytes, filetype="pdf")
        page = doc[0]
        cut_y = _find_title_cut_y(page)
        rect  = fitz.Rect(0, 0, page.rect.width, cut_y)
        return _render_region(page, rect, padding_pt=0)
    except Exception:
        return None


async def _analyze_figure_async(png_bytes: bytes, caption: str = "") -> dict:
    """
    Multimodal analysis: returns {"title": str, "description": str}.
    Caption (detected from PDF) is fed as a hint so the model can reproduce
    the original Fig/Table numbering in the title.
    """
    b64 = base64.standard_b64encode(png_bytes).decode()
    caption_hint = f"\n\nPDF内のキャプション（参考）：{caption}" if caption else ""

    prompt = (
        "この図表を分析し、以下の形式で日本語のみで出力してください。\n\n"
        "タイトル: [50字以内。PDF内にFig/Table番号があればそれを含める。"
        "例：Fig 1. カプラン・マイヤー生存曲線、Table 2. 患者背景の比較]\n"
        "解説: [150字以内。グラフの種類・表示データ・主要な傾向または統計的結果]"
        + caption_hint
    )

    try:
        resp = await async_client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=500,
            messages=[{
                "role": "user",
                "content": [
                    {"type": "image", "source": {
                        "type": "base64", "media_type": "image/png", "data": b64,
                    }},
                    {"type": "text", "text": prompt},
                ],
            }],
        )
        raw = resp.content[0].text

        title, desc = "", ""
        for line in raw.splitlines():
            if re.match(r'タイトル[:：]', line):
                title = re.sub(r'^タイトル[:：]\s*', '', line).strip()[:60]
            elif re.match(r'解説[:：]', line):
                desc = re.sub(r'^解説[:：]\s*', '', line).strip()

        # Fallback: if parsing failed, use the whole response as description
        if not title:
            title = (caption.split('.')[0] + '.').strip() if caption else "図表"
        if not desc:
            desc = raw.strip()[:300]

        return {"title": title, "description": desc}
    except Exception:
        return {"title": caption.split('\n')[0][:60] if caption else "図表", "description": "図表の解析に失敗しました。"}

# ── PPTX helpers ──────────────────────────────────────────────────────────────

def _strip_md(text: str) -> str:
    text = re.sub(r'\[\[(.+?)\]\]', r'\1', text)   # Obsidian links → plain text for PPTX
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    text = re.sub(r'^- ', '• ', text, flags=re.MULTILINE)
    return text.strip()


def _parse_sections(md: str) -> list[dict]:
    parts = re.split(r'(?m)^## ', md)
    sections = []
    for part in parts:
        if not part.strip():
            continue
        nl      = part.find('\n')
        title   = part[:nl].strip() if nl != -1 else part.strip()
        content = part[nl + 1:].strip() if nl != -1 else ''
        content = re.sub(r'^\s*---\s*$', '', content, flags=re.MULTILINE)
        content = re.sub(r'\n{3,}', '\n\n', content).strip()
        sections.append({'title': title, 'content': content})
    return sections


def _parse_md_table(content: str) -> list[list[str]]:
    rows = []
    for line in content.splitlines():
        line = line.strip()
        if not (line.startswith('|') and line.endswith('|')):
            continue
        if re.match(r'^\|[\s\-|:]+\|$', line):
            continue
        cells = [c.strip() for c in line[1:-1].split('|')]
        cells = [re.sub(r'\*\*(.+?)\*\*', r'\1', c) for c in cells]
        cells = [re.sub(r'\*(.+?)\*', r'\1', c) for c in cells]
        rows.append(cells)
    return rows


def _add_rect(slide, left, top, width, height, rgb: RGBColor):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = rgb
    s.line.fill.background()
    return s


def _add_title_bar(slide, slide_w, title: str, slide_num: int, total: int):
    bar = _add_rect(slide, 0, 0, slide_w, Inches(1.05), _NAVY)
    tf  = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = _WHITE

    nb = slide.shapes.add_textbox(slide_w - Inches(1.3), Inches(0.05), Inches(1.2), Inches(0.45))
    nt = nb.text_frame.paragraphs[0]
    nt.alignment = PP_ALIGN.RIGHT
    nr = nt.add_run()
    nr.text = f"{slide_num} / {total}"
    nr.font.size = Pt(11)
    nr.font.color.rgb = _WHITE


def _add_text_content(slide, content: str, top, slide_w, avail_h, font_pt: int = 15):
    pad = Inches(0.45)
    tb  = slide.shapes.add_textbox(pad, top, slide_w - pad * 2, avail_h)
    tf  = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    first = True
    for line in _strip_md(content).split('\n'):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = line
        r.font.size = Pt(font_pt)
        r.font.color.rgb = _BLACK


def _add_table_content(slide, rows: list[list[str]], top, slide_w, max_h):
    if not rows:
        return
    n_rows, n_cols = len(rows), len(rows[0])
    tbl_h = min(max_h, Inches(0.56) * n_rows)
    pad   = Inches(0.45)
    tbl_w = slide_w - pad * 2

    tbl = slide.shapes.add_table(n_rows, n_cols, pad, top, tbl_w, tbl_h).table
    if n_cols == 2:
        tbl.columns[0].width = int(tbl_w * 0.38)
        tbl.columns[1].width = int(tbl_w * 0.62)

    for r_idx, row in enumerate(rows):
        for c_idx, cell_text in enumerate(row[:n_cols]):
            cell = tbl.cell(r_idx, c_idx)
            cell.text_frame.word_wrap = True
            p   = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = cell_text
            run.font.size = Pt(14)
            run.font.color.rgb = _BLACK
            if r_idx == 0:
                run.font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = _LGRAY


def _add_section_content(slide, content: str, slide_w, slide_h):
    top     = Inches(1.2)
    avail_h = slide_h - top - Inches(0.15)
    has_tbl = bool(re.search(r'^\|.+\|$', content, re.MULTILINE))

    if has_tbl:
        rows      = _parse_md_table(content)
        non_table = re.sub(r'^\|.*\|$', '', content, flags=re.MULTILINE)
        non_table = re.sub(r'\n{2,}', '\n', non_table).strip()
        tbl_h     = min(Inches(4.8), Inches(0.56) * len(rows))
        _add_table_content(slide, rows, top, slide_w, tbl_h)
        if non_table:
            _add_text_content(slide, non_table, top + tbl_h + Inches(0.15), slide_w, avail_h - tbl_h, 14)
    else:
        _add_text_content(slide, content, top, slide_w, avail_h)


def _compute_fit(orig_w: int, orig_h: int, max_w: int, max_h: int) -> tuple[int, int]:
    if orig_w <= 0 or orig_h <= 0:
        return max_w, max_h
    scale = min(max_w / orig_w, max_h / orig_h)
    return int(orig_w * scale), int(orig_h * scale)


def _add_figure_slide(prs, slide, png_bytes: bytes, title: str, description: str,
                      slide_num: int, total: int):
    """Left 55%: image stretched to fill the full area. Right 43%: description."""
    w, h = prs.slide_width, prs.slide_height
    _add_rect(slide, 0, 0, w, h, _WHITE)
    _add_title_bar(slide, w, title, slide_num, total)

    content_top = Inches(1.15)
    content_h   = h - content_top - Inches(0.05)

    # ── Left 55 %: stretch image to fill every pixel of the left panel ───────
    img_area_w = int(w * 0.55)
    slide.shapes.add_picture(
        io.BytesIO(png_bytes),
        0, content_top,           # flush to left edge
        img_area_w, content_h,    # explicit width + height → fills area completely
    )

    # ── Right 43 %: description ───────────────────────────────────────────────
    txt_left = img_area_w + Inches(0.18)
    txt_w    = w - txt_left - Inches(0.1)
    tb = slide.shapes.add_textbox(txt_left, content_top + Inches(0.1), txt_w, content_h - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = description
    r.font.size = Pt(16)
    r.font.color.rgb = _BLACK


def _new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs, prs.slide_layouts[6]


def _add_title_slide(prs, blank, source_filename: str):
    w, h = prs.slide_width, prs.slide_height
    ts   = prs.slides.add_slide(blank)
    _add_rect(ts, 0, 0, w, h, _WHITE)

    tb = ts.shapes.add_textbox(Inches(1.5), Inches(2.2), w - Inches(3), Inches(2.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "抄読会資料"
    r.font.size = Pt(44)
    r.font.bold = True
    r.font.color.rgb = _NAVY

    if source_filename:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = source_filename.replace('.pdf', '')
        r2.font.size = Pt(20)
        r2.font.color.rgb = _MGRAY

    _add_rect(ts, Inches(1.5), Inches(4.7), w - Inches(3), Pt(3), _NAVY)


def _add_section_slides(prs, blank, sections: list[dict], total: int, start: int = 1):
    w, h = prs.slide_width, prs.slide_height
    for i, sec in enumerate(sections, start):
        sl = prs.slides.add_slide(blank)
        _add_rect(sl, 0, 0, w, h, _WHITE)
        _add_title_bar(sl, w, sec['title'], i, total)
        _add_section_content(sl, sec['content'], w, h)


def _add_title_screenshot_slide(prs, slide, png_bytes: bytes, slide_num: int, total: int):
    """Full-page screenshot of the paper's title / authors / institution area."""
    w, h = prs.slide_width, prs.slide_height
    _add_rect(slide, 0, 0, w, h, _WHITE)
    _add_title_bar(slide, w, "タイトル・著者・研究機関", slide_num, total)

    content_top = Inches(1.15)
    content_h   = h - content_top - Inches(0.1)
    content_w   = w - Inches(0.5)

    img = PILImage.open(io.BytesIO(png_bytes))
    img_w, img_h = img.size
    if img_h == 0:
        return
    if img_w / img_h > content_w / content_h:
        disp_w = int(content_w)
        disp_h = int(content_w * img_h / img_w)
    else:
        disp_h = int(content_h)
        disp_w = int(content_h * img_w / img_h)
    left = (w - disp_w) // 2
    top  = content_top + (content_h - disp_h) // 2
    slide.shapes.add_picture(io.BytesIO(png_bytes), left, top, disp_w, disp_h)


def generate_pptx(markdown: str, source_filename: str = "") -> bytes:
    sections = _parse_sections(markdown)
    total    = 1 + len(sections)
    prs, blank = _new_prs()
    _add_title_slide(prs, blank, source_filename)
    _add_section_slides(prs, blank, sections, total)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def generate_pptx_with_figures(
    markdown:         str,
    figures:          list[dict],
    source_filename:  str          = "",
    title_screenshot: bytes | None = None,
) -> bytes:
    sections     = _parse_sections(markdown)
    n_sec, n_fig = len(sections), len(figures)
    n_ts         = 1 if title_screenshot else 0
    total        = 1 + n_ts + n_sec + n_fig
    prs, blank   = _new_prs()

    _add_title_slide(prs, blank, source_filename)

    if title_screenshot:
        sl = prs.slides.add_slide(blank)
        _add_title_screenshot_slide(prs, sl, title_screenshot, 1, total)

    _add_section_slides(prs, blank, sections, total, start=1 + n_ts)

    for j, fig in enumerate(figures, 1):
        sl = prs.slides.add_slide(blank)
        _add_figure_slide(
            prs, sl,
            fig["png_bytes"],
            fig["title"],
            fig["description"],
            n_ts + n_sec + j, total,
        )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _add_progress_title_slide(prs, blank):
    w, h = prs.slide_width, prs.slide_height
    ts   = prs.slides.add_slide(blank)
    _add_rect(ts, 0, 0, w, h, _WHITE)
    _add_rect(ts, 0, 0, w, Inches(0.5), _NAVY)
    _add_rect(ts, 0, h - Inches(0.5), w, Inches(0.5), _NAVY)

    today = datetime.date.today().strftime("%Y年%m月%d日")

    tb = ts.shapes.add_textbox(Inches(1.5), Inches(1.5), w - Inches(3), Inches(3.2))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "医学論文 抄読会資料ジェネレーター"
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = _NAVY

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "開発進捗報告"
    r2.font.size = Pt(38)
    r2.font.bold = True
    r2.font.color.rgb = _NAVY

    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = today
    r3.font.size = Pt(16)
    r3.font.color.rgb = _MGRAY

    _add_rect(ts, Inches(1.5), Inches(5.0), w - Inches(3), Pt(3), _NAVY)


def generate_progress_pptx(dynamic_markdown: str) -> bytes:
    dynamic_sections = _parse_sections(dynamic_markdown)
    all_sections     = FIXED_PROJECT_SECTIONS + dynamic_sections
    total            = 1 + len(all_sections)
    prs, blank       = _new_prs()
    _add_progress_title_slide(prs, blank)
    _add_section_slides(prs, blank, all_sections, total)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _write_obsidian_note(subfolder: str, filename: str, content: str) -> str:
    vault  = pathlib.Path(OBSIDIAN_VAULT_PATH)
    folder = vault / subfolder
    folder.mkdir(parents=True, exist_ok=True)
    note_path = folder / filename
    note_path.write_text(content, encoding="utf-8")
    return str(note_path)

# ── API endpoints ──────────────────────────────────────────────────────────────

@app.get("/health")
def health():
    return {"status": "ok", "message": "バックエンド正常稼働中"}


@app.post("/estimate-cost")
async def estimate_cost(file: UploadFile = File(...)):
    """
    Estimate the number of tokens and API cost (JPY) before running full processing.

    Counts:
      - Text tokens  : len(extracted_text) * 1.2  (rough approximation)
      - Image tokens : (n_figures + 1 title screenshot) * TOKENS_PER_IMAGE
    Output ceiling:
      - MAX_TOKENS_SUMMARIZE  (1 call)
      - MAX_TOKENS_FIGURE * n_figures (parallel vision calls)
    """
    pdf_bytes = await file.read()

    # ── Page count ────────────────────────────────────────────────────────────
    doc         = fitz.open(stream=pdf_bytes, filetype="pdf")
    total_pages = len(doc)
    doc.close()

    # ── Text tokens ───────────────────────────────────────────────────────────
    try:
        paper_text  = extract_text(pdf_bytes)
        text_tokens = int(len(paper_text) * 1.2)
    except Exception:
        text_tokens = 0

    # ── Figure / Table count ──────────────────────────────────────────────────
    try:
        figures          = extract_figures(pdf_bytes)
        detected_figures = len(figures)
    except Exception:
        detected_figures = 0

    # ── Title screenshot (1 image if extractable) ────────────────────────────
    title_ss    = extract_title_screenshot(pdf_bytes)
    n_images    = detected_figures + (1 if title_ss else 0)

    # ── Token totals ──────────────────────────────────────────────────────────
    estimated_input_tokens  = text_tokens + n_images * TOKENS_PER_IMAGE
    estimated_output_tokens = MAX_TOKENS_SUMMARIZE + detected_figures * MAX_TOKENS_FIGURE

    # ── Cost (ceiling in JPY) ─────────────────────────────────────────────────
    cost_usd = (
        estimated_input_tokens  * PRICE_INPUT_TEXT_USD / 1_000_000
        + estimated_output_tokens * PRICE_OUTPUT_USD      / 1_000_000
    )
    estimated_cost_jpy = math.ceil(cost_usd * USD_TO_JPY)

    return {
        "total_pages":             total_pages,
        "detected_figures":        detected_figures,
        "estimated_input_tokens":  estimated_input_tokens,
        "estimated_output_tokens": estimated_output_tokens,
        "estimated_cost_jpy":      estimated_cost_jpy,
    }


@app.post("/summarize")
async def summarize(file: UploadFile = File(...)):
    if not file.filename.lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="PDFファイルのみ対応しています")
    pdf_bytes = await file.read()
    try:
        paper_text = extract_text(pdf_bytes)
    except Exception as e:
        raise HTTPException(status_code=422, detail=f"PDF読み込みエラー: {e}")
    if not paper_text.strip():
        raise HTTPException(status_code=422, detail="PDFからテキストを抽出できませんでした（スキャンPDFは非対応）")
    try:
        response = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=4096,
            system=SYSTEM_PROMPT,
            messages=[{"role": "user", "content": build_prompt(paper_text)}],
        )
        summary = response.content[0].text
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"AI生成エラー: {e}")
    return {"filename": file.filename, "summary": summary}


class ExportRequest(BaseModel):
    summary:  str
    filename: str = "論文"


class ObsidianNoteRequest(BaseModel):
    summary:  str
    filename: str = "論文"


@app.post("/export/pptx")
async def export_pptx(req: ExportRequest):
    try:
        pptx_bytes = generate_pptx(req.summary, req.filename)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PPTX生成エラー: {e}")
    stem = req.filename.replace('.pdf', '')
    dl   = f"{stem}_抄読会資料.pptx"
    return StreamingResponse(
        io.BytesIO(pptx_bytes),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename*=UTF-8''{quote(dl)}"},
    )


@app.post("/export/pptx-with-figures")
async def export_pptx_with_figures(
    file:     UploadFile = File(...),
    summary:  str        = Form(...),
    filename: str        = Form("論文"),
):
    pdf_bytes = await file.read()

    title_ss = extract_title_screenshot(pdf_bytes)

    try:
        raw_figures = extract_figures(pdf_bytes)
    except Exception as e:
        raise HTTPException(status_code=422, detail=f"図表抽出エラー: {e}")

    if not raw_figures and title_ss is None:
        raise HTTPException(
            status_code=422,
            detail="PDFからタイトルページ・図表を抽出できませんでした",
        )

    analyses = await asyncio.gather(
        *[_analyze_figure_async(f["png_bytes"], f["caption"]) for f in raw_figures]
    )
    analyzed = [{**fig, **ana} for fig, ana in zip(raw_figures, analyses)]

    try:
        pptx_bytes = generate_pptx_with_figures(summary, analyzed, filename, title_ss)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PPTX生成エラー: {e}")

    stem = filename.replace('.pdf', '')
    dl   = f"{stem}_抄読会資料_図表付.pptx"
    return StreamingResponse(
        io.BytesIO(pptx_bytes),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename*=UTF-8''{quote(dl)}"},
    )


@app.post("/export/obsidian-note")
async def export_obsidian_note(req: ObsidianNoteRequest):
    if not OBSIDIAN_VAULT_PATH:
        raise HTTPException(
            status_code=503,
            detail="OBSIDIAN_VAULT_PATH が未設定です。.env の OBSIDIAN_VAULT_HOST_PATH を確認してください。",
        )
    today = datetime.date.today().isoformat()
    stem  = req.filename.replace(".pdf", "")
    note  = (
        f"---\ntitle: {stem}\ndate: {today}\ntags:\n"
        f"  - 医学論文\n  - 抄読会\nsource: {req.filename}\n---\n\n"
        f"[[論文要約AIプロジェクトMOC]]\n\n{req.summary}\n"
    )
    try:
        path = _write_obsidian_note("論文要約", f"{stem}.md", note)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Obsidianへの書き込みエラー: {e}")
    return {"status": "ok", "note_name": f"{stem}.md", "path": path}


@app.post("/sync/devlog")
async def sync_devlog():
    if not OBSIDIAN_VAULT_PATH:
        raise HTTPException(
            status_code=503,
            detail="OBSIDIAN_VAULT_PATH が未設定です。.env の OBSIDIAN_VAULT_HOST_PATH を確認してください。",
        )
    if not DEV_LOG_PATH.exists():
        raise HTTPException(status_code=404, detail="dev_log.md が見つかりません")
    dev_log = DEV_LOG_PATH.read_text(encoding="utf-8")
    today   = datetime.date.today().isoformat()
    note    = (
        f"---\ntitle: Paper Summarizer AI 開発ログ\ntags:\n"
        f"  - 開発ログ\n  - paper-summarizer-ai\nupdated: {today}\n---\n\n"
        f"[[論文要約AIプロジェクトMOC]]\n\n{dev_log}\n"
    )
    try:
        path = _write_obsidian_note("開発ログ", "paper-summarizer-ai.md", note)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Obsidianへの書き込みエラー: {e}")
    return {"status": "ok", "path": path}


@app.post("/export/progress-pptx")
async def export_progress_pptx():
    dev_log    = DEV_LOG_PATH.read_text(encoding="utf-8") if DEV_LOG_PATH.exists() else ""
    dynamic_md = ""
    if dev_log.strip():
        try:
            resp = client.messages.create(
                model="claude-sonnet-4-6",
                max_tokens=2048,
                system=PROGRESS_SYSTEM_PROMPT,
                messages=[{"role": "user", "content": build_progress_prompt(dev_log)}],
            )
            dynamic_md = resp.content[0].text
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"AI生成エラー: {e}")
    try:
        pptx_bytes = generate_progress_pptx(dynamic_md)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PPTX生成エラー: {e}")
    today = datetime.date.today().strftime("%Y%m%d")
    dl    = f"進捗報告_{today}.pptx"
    return StreamingResponse(
        io.BytesIO(pptx_bytes),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename*=UTF-8''{quote(dl)}"},
    )
